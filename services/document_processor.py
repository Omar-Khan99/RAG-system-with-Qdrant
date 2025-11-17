# src/services/document_processor.py
from ollama import Client
from docx import Document
import PyPDF2
import pandas as pd
from pptx import Presentation
import os
from typing import Dict, Any, List
import re

class DocumentProcessor:
    def __init__(self):
        pass

    
    def extract_text_from_docx(self, file_path: str) -> str:
        """استخراج النص من ملفات DOCX"""
        try:
            doc = Document(file_path)
            text = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text.append(paragraph.text)
            return '\n'.join(text)
        except Exception as e:
            raise Exception(f"Error reading DOCX file: {str(e)}")
    
    def extract_text_from_pdf(self, file_path: str) -> str:
        """استخراج النص من ملفات PDF"""
        try:
            text = []
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text.append(page_text)
            return '\n'.join(text)
        except Exception as e:
            raise Exception(f"Error reading PDF file: {str(e)}")
    
    def extract_text_from_txt(self, file_path: str) -> str:
        """استخراج النص من ملفات TXT"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except Exception as e:
            raise Exception(f"Error reading TXT file: {str(e)}")
    
    def extract_text_from_pptx(self, file_path: str) -> str:
        """استخراج النص من ملفات PPTX"""
        try:
            prs = Presentation(file_path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text.append(shape.text)
            return '\n'.join(text)
        except Exception as e:
            raise Exception(f"Error reading PPTX file: {str(e)}")
    
    def extract_text_from_csv(self, file_path: str) -> str:
        """استخراج النص من ملفات CSV"""
        try:
            df = pd.read_csv(file_path)
            text = ["Columns: " + ", ".join(df.columns.tolist())]
            
            for idx, row in df.head(10).iterrows():
                text.append(f"Row {idx}: " + ", ".join([str(x) for x in row.values]))
            
            return '\n'.join(text)
        except Exception as e:
            raise Exception(f"Error reading CSV file: {str(e)}")
    
    def extract_text_from_excel(self, file_path: str) -> str:
        """استخراج النص من ملفات Excel"""
        try:
            df = pd.read_excel(file_path)
            text = ["Columns: " + ", ".join(df.columns.tolist())]
            
            for idx, row in df.head(10).iterrows():
                text.append(f"Row {idx}: " + ", ".join([str(x) for x in row.values]))
            
            return '\n'.join(text)
        except Exception as e:
            raise Exception(f"Error reading Excel file: {str(e)}")
    
    def chunk_text(self, text: str, chunk_size: int = 500, overlap: int = 50) -> List[Dict[str, Any]]:
        """تقسيم النص إلى أجزاء متداخلة"""
        sentences = re.split(r'[.!?]+', text)
        current_chunk = []
        current_length = 0
        chunks = []

        for sentence in sentences:
            sentence = sentence.strip()
            if not sentence:
                continue

            sentence_length = len(sentence.split())

            if sentence_length + current_length <= chunk_size:
                current_chunk.append(sentence)
                current_length += sentence_length
            else:
                if current_chunk:
                    chunks.append(' '.join(current_chunk))

                if overlap > 0 and chunks:
                    last_word_chunk = chunks[-1].split()[-overlap:]
                    current_chunk = last_word_chunk + [sentence]
                    current_length = len(current_chunk)
                else:
                    current_chunk = [sentence]
                    current_length = sentence_length

        if current_chunk:
            chunks.append(' '.join(current_chunk))

        return [{'text': chunk, 'chunk_id': i, 'total_words': len(chunk.split())} 
                for i, chunk in enumerate(chunks)]
    
    def extract_metadata(self, file_path: str, text: str) -> Dict[str, Any]:
        """استخراج البيانات الوصفية من المستند"""
        file_name = os.path.basename(file_path)
        file_size = os.path.getsize(file_path)
        file_extension = os.path.splitext(file_path)[1].lower()
        
        lines = text.split('\n')
        first_lines = lines[:5]
        
        return {
            'file_name': file_name,
            'file_size': file_size,
            'file_type': file_extension,
            'processed_date': pd.Timestamp.now().isoformat(),
            'total_chars': len(text),
            'total_words': len(text.split()),
            'first_lines': first_lines,
        }
    
    def process_document(self, file_path: str) -> str:
        """معالجة المستند بناءً على امتداد الملف"""
        ext = os.path.splitext(file_path)[1].lower()
        
        processors = {
            '.docx': self.extract_text_from_docx,
            '.pdf': self.extract_text_from_pdf,
            '.txt': self.extract_text_from_txt,
            '.pptx': self.extract_text_from_pptx,
            '.csv': self.extract_text_from_csv,
            '.xlsx': self.extract_text_from_excel,
            '.xls': self.extract_text_from_excel
        }
        
        if ext not in processors:
            raise ValueError(f"Unsupported file type: {ext}")
        
        return processors[ext](file_path)